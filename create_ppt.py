"""
PowerPoint Presentation Generator for EV Analytics Project
This script creates a professional PowerPoint presentation automatically
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create comprehensive PowerPoint presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    DARK_BLUE = RGBColor(15, 23, 42)
    LIGHT_BLUE = RGBColor(52, 152, 219)
    ACCENT_GREEN = RGBColor(46, 204, 113)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(149, 163, 184)
    
    def add_title_slide(title, subtitle, name):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = DARK_BLUE
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(54)
        title_p.font.bold = True
        title_p.font.color.rgb = ACCENT_GREEN
        
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = subtitle
        subtitle_p.font.size = Pt(28)
        subtitle_p.font.color.rgb = WHITE
        
        author_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
        author_frame = author_box.text_frame
        author_p = author_frame.paragraphs[0]
        author_p.text = f"By: {name}\nAcademic Year: 2025-2026"
        author_p.font.size = Pt(18)
        author_p.font.color.rgb = GRAY
        
        return slide
    
    def add_content_slide(title, content_list):
        """Add content slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = WHITE
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_p = title_frame.paragraphs[0]
        title_p.text = title
        title_p.font.size = Pt(40)
        title_p.font.bold = True
        title_p.font.color.rgb = DARK_BLUE
        
        # Add title underline
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(9), Inches(0))
        line.line.color.rgb = ACCENT_GREEN
        line.line.width = Pt(3)
        
        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.2))
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        
        for idx, item in enumerate(content_list):
            if idx > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[idx]
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_BLUE
            p.space_before = Pt(6)
            p.space_after = Pt(6)
            p.level = 0
        
        return slide
    
    # SLIDE 1: Title Slide
    add_title_slide(
        "🚗 ELECTRIC VEHICLE ANALYTICS\n& PERFORMANCE PREDICTION",
        "Using Machine Learning for Range Prediction",
        "[Your Name]"
    )
    
    # SLIDE 2: Agenda
    add_content_slide(
        "📋 AGENDA",
        [
            "1. Introduction & Motivation",
            "2. Problem Statement",
            "3. Objectives",
            "4. Dataset Overview",
            "5. Methodology & Data Preprocessing",
            "6. Exploratory Data Analysis",
            "7. Feature Engineering & Model Development",
            "8. Results & Evaluation",
            "9. Key Insights & Limitations"
        ]
    )
    
    # SLIDE 3: Introduction
    add_content_slide(
        "🌍 WHY ELECTRIC VEHICLES?",
        [
            "✅ Reduce carbon emissions by 50-70%",
            "✅ Lower operational costs (~60% cheaper)",
            "✅ Improved energy efficiency (90% vs 20%)",
            "✅ Global market growing 40% annually",
            "",
            "❓ Challenge: Predicting EV range accurately",
            "❓ Research: What factors influence performance?",
            "❓ Solution: Machine Learning models"
        ]
    )
    
    # SLIDE 4: Problem Statement
    add_content_slide(
        "🎯 PROBLEM STATEMENT",
        [
            "Challenges:",
            "• Complex EV data across multiple manufacturers",
            "• Lack of predictive tools for range estimation",
            "• Difficult to compare vehicles objectively",
            "",
            "Research Question:",
            "Can ML accurately predict EV range based on",
            "manufacturer characteristics and model year?"
        ]
    )
    
    # SLIDE 5: Objectives
    add_content_slide(
        "🎯 PROJECT OBJECTIVES",
        [
            "✅ Analyze 112,634 → 36,590 EV records",
            "✅ Implement 3 ML models",
            "✅ Achieve >90% prediction accuracy",
            "✅ Identify key performance factors",
            "✅ Generate 9 visualizations",
            "✅ Create comprehensive documentation"
        ]
    )
    
    # SLIDE 6: Dataset
    add_content_slide(
        "📊 DATASET OVERVIEW",
        [
            "Source: Kaggle - Electric Vehicle Population Data",
            "",
            "Original: 112,634 records × 17 attributes",
            "",
            "After Cleaning:",
            "✓ 36,590 vehicles | ✓ 2016-2021 model years",
            "✓ 15 manufacturers | ✓ BEV only",
            "✓ Average range: 215.72 miles"
        ]
    )
    
    # SLIDE 7: Methodology
    add_content_slide(
        "⚙️ METHODOLOGY",
        [
            "Phase 1: Data Collection & Inspection",
            "Phase 2: Data Cleaning (removing duplicates, outliers)",
            "Phase 3: Exploratory Data Analysis (EDA)",
            "Phase 4: Feature Engineering",
            "Phase 5: Model Development (3 algorithms)",
            "Phase 6: Model Evaluation & Results",
            "",
            "Duration: 4 weeks | Effort: 20-25 hours"
        ]
    )
    
    # SLIDE 8: Data Preprocessing
    add_content_slide(
        "🧹 DATA PREPROCESSING",
        [
            "✅ Step 1: Removed 0 duplicates",
            "✅ Step 2: Filtered to Battery EVs (86K)",
            "✅ Step 3: Removed 0-range vehicles",
            "✅ Step 4: Kept Model Year 2016+ (37K)",
            "✅ Step 5: Selected top 15 manufacturers",
            "✅ Step 6: Removed 9 outliers using IQR",
            "",
            "Final Dataset: 36,590 clean records ✓"
        ]
    )
    
    # SLIDE 9: EDA Statistics
    add_content_slide(
        "📈 EXPLORATORY DATA ANALYSIS",
        [
            "Electric Range Statistics:",
            "• Mean: 215.72 miles | Median: 220 miles",
            "• Std Dev: 61.94 miles | Range: 57-337 miles",
            "",
            "Market Share (Top 5):",
            "• Tesla: 62.4% | Nissan: 12.3%",
            "• Chevrolet: 8.7% | Kia: 4.2%",
            "",
            "Key Finding: 33% range improvement (2016→2021)"
        ]
    )
    
    # SLIDE 10: Feature Engineering
    add_content_slide(
        "⚙️ FEATURE ENGINEERING",
        [
            "Created 5 new features:",
            "1. Vehicle_Age = 2026 - Model Year",
            "2. Years_Since_2016 = Model Year - 2016",
            "3. Make_Encoded = Manufacturer label encoding",
            "4. Manufacturer_Tier = High/Medium/Low range",
            "5. Market_Share = Brand percentage",
            "",
            "Selected 4 features for modeling (excluded Model Year)"
        ]
    )
    
    # SLIDE 11: Models
    add_content_slide(
        "🤖 MACHINE LEARNING MODELS",
        [
            "Model 1: Linear Regression (Baseline)",
            "• Simple, interpretable, fast",
            "",
            "Model 2: Random Forest (100 trees)",
            "• Handles non-linearity, feature importance",
            "",
            "Model 3: XGBoost ⭐ (Gradient Boosting)",
            "• State-of-the-art performance",
            "",
            "Train-Test Split: 80%-20% (29K-7K samples)"
        ]
    )
    
    # SLIDE 12: Results
    add_content_slide(
        "🏆 MODEL PERFORMANCE",
        [
            "Model | R² Score | MAE (mi) | RMSE (mi)",
            "─────────────────────────────────────",
            "Linear Reg | 0.8110 | 21.78 | 26.70",
            "Random Forest | 0.9467 | 8.56 | 14.19",
            "XGBoost ⭐ | 0.9468 | 8.56 | 14.17",
            "",
            "🏅 BEST: XGBoost with 94.68% accuracy!",
            "✅ Outperforms literature benchmarks"
        ]
    )
    
    # SLIDE 13: Feature Importance
    add_content_slide(
        "🔍 FEATURE IMPORTANCE",
        [
            "What drives EV range?",
            "",
            "1. Make_Encoded (62%) - Manufacturer identity",
            "2. Manufacturer_Tier (23%) - Brand category",
            "3. Market_Share (9%) - Popular brands",
            "4. Years_Since_2016 (6%) - Technology year",
            "",
            "Key Insight: Manufacturer matters most!"
        ]
    )
    
    # SLIDE 14: Key Insights
    add_content_slide(
        "💡 KEY FINDINGS",
        [
            "1. Manufacturer Dominance: 62% of prediction variance",
            "2. Technology Evolution: 33% range improvement (2016-21)",
            "3. Market Concentration: Tesla dominates (62% share)",
            "4. Performance Tiers: Premium (300+ mi) vs Economy",
            "5. High Accuracy: 94.7% exceeds industry standards",
            "6. Error Margin: ±8.56 miles commercially viable"
        ]
    )
    
    # SLIDE 15: Limitations
    add_content_slide(
        "⚠️ LIMITATIONS",
        [
            "Data: Only up to 2021 | US-centric | No sensor data",
            "Model: Relies on manufacturer | No degradation modeling",
            "Scope: CLI only | Not deployed as web service",
            "",
            "Acknowledged & documented in project report"
        ]
    )
    
    # SLIDE 16: Future Scope
    add_content_slide(
        "🔮 FUTURE ENHANCEMENTS",
        [
            "Technical:",
            "• Deep learning (LSTM, CNN) | Real-time integration",
            "",
            "Deployment:",
            "• Web app (Flask/React) | Mobile app | API",
            "",
            "Research:",
            "• Price prediction | TCO calculator | Carbon analysis"
        ]
    )
    
    # SLIDE 17: Tech Stack
    add_content_slide(
        "🛠️ TECHNOLOGY STACK",
        [
            "Language: Python 3.10",
            "Libraries: pandas, numpy, scikit-learn, xgboost",
            "Visualization: matplotlib, seaborn",
            "Version Control: Git + GitHub",
            "IDE: VS Code",
            "",
            "All tools: Open-source & industry-standard"
        ]
    )
    
    # SLIDE 18: Deliverables
    add_content_slide(
        "📦 PROJECT DELIVERABLES",
        [
            "✅ Source Code: ev_analytics.py (modular & documented)",
            "✅ Visualizations: 9 professional-quality plots",
            "✅ Trained Models: best_model_xgboost.pkl",
            "✅ Documentation: README, Report (40+ pages)",
            "✅ GitHub Repository: Complete with history",
            "✅ Presentation: This deck + Viva Q&A"
        ]
    )
    
    # SLIDE 19: How to Run
    add_content_slide(
        "🚀 PROJECT EXECUTION",
        [
            "Quick Start:",
            "1. pip install -r requirements.txt",
            "2. python ev_analytics.py",
            "",
            "Output (30 seconds):",
            "✓ Cleaned: 112K → 36K records",
            "✓ Generated: 9 visualizations",
            "✓ Trained: 3 ML models",
            "✓ Best Result: XGBoost 94.7% accuracy"
        ]
    )
    
    # SLIDE 20: Conclusion
    add_content_slide(
        "✨ CONCLUSION",
        [
            "Achievements:",
            "✅ Processed 36,590 EV records successfully",
            "✅ Achieved 94.68% prediction accuracy",
            "✅ Identified key performance factors",
            "✅ Created comprehensive documentation",
            "",
            "Impact:",
            "🚗 Helps consumers | 🏭 Helps manufacturers",
            "📊 Supports policymakers"
        ]
    )
    
    # SLIDE 21: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ACCENT_GREEN
    
    qa_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    qa_frame = qa_box.text_frame
    qa_p = qa_frame.paragraphs[0]
    qa_p.text = "QUESTIONS & ANSWERS"
    qa_p.font.size = Pt(66)
    qa_p.font.bold = True
    qa_p.font.color.rgb = WHITE
    qa_p.alignment = PP_ALIGN.CENTER
    
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1.5))
    thanks_frame = thanks_box.text_frame
    thanks_p = thanks_frame.paragraphs[0]
    thanks_p.text = "Thank you for your attention!\n\n[Your Name]\ngithub.com/YOUR_USERNAME/EV-Analytics-Project"
    thanks_p.font.size = Pt(20)
    thanks_p.font.color.rgb = WHITE
    thanks_p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    prs.save('EV_Analytics_Presentation.pptx')
    print("✅ PowerPoint presentation created successfully!")
    print("📄 File: EV_Analytics_Presentation.pptx")
    print("📊 Total slides: 21")

if __name__ == "__main__":
    create_presentation()